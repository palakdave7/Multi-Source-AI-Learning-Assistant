import os
import re
import httpx
from bs4 import BeautifulSoup
from utils.text import clean_text, chunk_text
from pypdf import PdfReader
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi
from typing import List


# ---------- helpers ----------

def chunk_text(text: str, chunk_size: int = 1500, overlap: int = 150) -> List[str]:
    """Split text into overlapping chunks by word count."""
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunks.append(" ".join(words[start:end]))
        start += chunk_size - overlap
    return [c for c in chunks if c.strip()]


# ---------- YouTube ----------

def _extract_youtube_id(url: str) -> str:
    patterns = [
        r"(?:v=|\/)([0-9A-Za-z_-]{11})",
        r"youtu\.be\/([0-9A-Za-z_-]{11})",
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    raise ValueError("Could not extract YouTube video ID from URL")


def process_youtube(url: str) -> dict:
    video_id = _extract_youtube_id(url)
    try:
        api = YouTubeTranscriptApi()
        transcript_list = api.fetch(video_id)
    except Exception:
        raise ValueError("Could not fetch transcript. The video may have no captions or be private.")

    chunks = []
    buffer_text = []
    buffer_start = 0.0
    buffer_words = 0

    for entry in transcript_list:
        words = entry.text.split()
        if buffer_words == 0:
            buffer_start = entry.start
        buffer_text.append(entry.text)
        buffer_words += len(words)
        if buffer_words >= 150:
            timestamp = _format_timestamp(buffer_start)
            chunks.append({
                "text": " ".join(buffer_text),
                "ref": f"at {timestamp} in the video",
            })
            buffer_text = []
            buffer_words = 0

    if buffer_text:
        chunks.append({
            "text": " ".join(buffer_text),
            "ref": f"at {_format_timestamp(buffer_start)} in the video",
        })

    full_text = " ".join(e.text for e in transcript_list)
    summary_snippet = full_text[:500] + "..." if len(full_text) > 500 else full_text

    return {
        "type": "youtube",
        "label": f"YouTube: {video_id}",
        "url": url,
        "summary_snippet": summary_snippet,
        "chunks": chunks,
    }

def _format_timestamp(seconds: float) -> str:
    m, s = divmod(int(seconds), 60)
    h, m = divmod(m, 60)
    if h:
        return f"{h}:{m:02d}:{s:02d}"
    return f"{m}:{s:02d}"


# ---------- PDF ----------

def process_pdf(file_bytes: bytes, filename: str) -> dict:
    from io import BytesIO
    reader = PdfReader(BytesIO(file_bytes))
    chunks = []

    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if not text.strip():
            continue
        for chunk in chunk_text(text):
            chunks.append({
                "text": chunk,
                "ref": f"page {page_num} of {filename}",
            })

    full_text = " ".join(c["text"] for c in chunks)
    summary_snippet = full_text[:500] + "..." if len(full_text) > 500 else full_text

    return {
        "type": "pdf",
        "label": f"PDF: {filename}",
        "summary_snippet": summary_snippet,
        "chunks": chunks,
    }


# ---------- PPTX ----------

def process_pptx(file_bytes: bytes, filename: str) -> dict:
    from io import BytesIO
    prs = Presentation(BytesIO(file_bytes))
    chunks = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slide_text = " ".join(texts)
        if not slide_text.strip():
            continue
        for chunk in chunk_text(slide_text):
            chunks.append({
                "text": chunk,
                "ref": f"slide {slide_num} of {filename}",
            })

    full_text = " ".join(c["text"] for c in chunks)
    summary_snippet = full_text[:500] + "..." if len(full_text) > 500 else full_text

    return {
        "type": "pptx",
        "label": f"PPTX: {filename}",
        "summary_snippet": summary_snippet,
        "chunks": chunks,
    }


# ---------- Webpage ----------

def process_url(url: str) -> dict:
    response = httpx.get(url, timeout=15, follow_redirects=True, headers={
        "User-Agent": "Mozilla/5.0 (compatible; SamasocialBot/1.0)"
    })
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")

    # Remove nav, footer, scripts
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    title = soup.title.string.strip() if soup.title else url
    text = soup.get_text(separator=" ", strip=True)
    text = re.sub(r"\s+", " ", text)

    chunks = []
    for chunk in chunk_text(text):
        chunks.append({
            "text": chunk,
            "ref": f"from {url}",
        })

    summary_snippet = text[:500] + "..." if len(text) > 500 else text

    return {
        "type": "url",
        "label": f"Web: {title[:60]}",
        "url": url,
        "summary_snippet": summary_snippet,
        "chunks": chunks,
    }